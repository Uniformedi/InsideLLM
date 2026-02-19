"""
title: DLP Filter Pipeline
author: Uniformedi LLC
version: 2.0.0
description: Data Loss Prevention filter that scans outbound messages AND uploaded
             files (Excel, CSV, PDF, Word, PowerPoint, text) for PII, PHI, credit
             card numbers, SSNs, API keys, and connection strings before they reach
             the Claude API. Blocks or redacts sensitive data.
"""

import re
import os
import json
import logging
import glob as glob_mod
from typing import Optional, Dict, List, Tuple
from pydantic import BaseModel, Field

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("dlp-pipeline")

UPLOAD_DIR = os.environ.get("UPLOAD_DIR", "/app/backend/data/uploads")


class Pipeline:
    """
    Open WebUI Filter Pipeline for Data Loss Prevention.

    This runs as a pre-processing filter on every user message before it is
    sent to the LLM API. It scans for sensitive data patterns in both message
    text and uploaded files, then either blocks the message entirely or redacts
    the sensitive portions.

    Supported file formats: Excel (.xlsx/.xls), CSV, TSV, PDF, Word (.docx),
    PowerPoint (.pptx), and plain text files (.txt, .md, .json, .xml, .yaml,
    .py, .js, .log).
    """

    class Valves(BaseModel):
        """Configurable settings exposed in Open WebUI admin panel."""
        enabled: bool = Field(
            default=True,
            description="Enable/disable the DLP filter"
        )
        mode: str = Field(
            default="block",
            description="Action mode: 'block' (reject message) or 'redact' (replace matches)"
        )
        block_ssn: bool = Field(
            default=True,
            description="Detect Social Security Numbers"
        )
        block_credit_cards: bool = Field(
            default=True,
            description="Detect credit card numbers"
        )
        block_phi: bool = Field(
            default=True,
            description="Detect Protected Health Information patterns"
        )
        block_credentials: bool = Field(
            default=True,
            description="Detect API keys, passwords, connection strings"
        )
        block_bank_accounts: bool = Field(
            default=True,
            description="Detect bank account and routing numbers"
        )
        block_standalone_dates: bool = Field(
            default=True,
            description="Detect standalone date patterns (MM/DD/YYYY, YYYY-MM-DD) that may be dates of birth"
        )
        scan_file_uploads: bool = Field(
            default=True,
            description="Scan uploaded files (Excel, CSV, PDF, Word, etc.) for sensitive data"
        )
        max_file_size_mb: int = Field(
            default=50,
            description="Maximum file size in MB to scan (larger files are skipped)"
        )
        log_detections: bool = Field(
            default=True,
            description="Log detected patterns (without the actual data)"
        )
        custom_patterns: str = Field(
            default="{}",
            description='JSON object of custom patterns: {"name": "regex"}'
        )

    def __init__(self):
        self.name = "DLP Filter"
        self.file_handler = False
        self.valves = self.Valves()

        # Core detection patterns
        self.patterns: Dict[str, Dict] = {
            "ssn": {
                "regex": r'\b\d{3}[-\u2013\u2014\s]?\d{2}[-\u2013\u2014\s]?\d{4}\b',
                "description": "Social Security Number",
                "valve": "block_ssn",
                "severity": "critical",
            },
            "ssn_labeled": {
                "regex": r'(?:social\s*security(?:\s*(?:number|no\.?|num\.?|#))?|ssn|ss\s*#|ss\s*no\.?)[\s:.=#]*\d{3}[-\u2013\u2014\s]?\d{2}[-\u2013\u2014\s]?\d{4}',
                "description": "Social Security Number (labeled)",
                "valve": "block_ssn",
                "severity": "critical",
            },
            "credit_card": {
                "regex": r'\b(?:4\d{3}|5[1-5]\d{2}|3[47]\d{2}|6(?:011|5\d{2}))[-\s]?\d{4}[-\s]?\d{4}[-\s]?\d{4}\b',
                "description": "Credit Card Number",
                "valve": "block_credit_cards",
                "severity": "critical",
            },
            "credit_card_generic": {
                "regex": r'\b(?:\d{4}[-\s]){3}\d{4}\b',
                "description": "Potential Credit Card Number",
                "valve": "block_credit_cards",
                "severity": "high",
            },
            "phi_mrn": {
                "regex": r'\b(?:MRN|Medical Record|Patient ID)[\s:#]*\d{5,}\b',
                "description": "Medical Record Number",
                "valve": "block_phi",
                "severity": "critical",
            },
            "phi_dob": {
                "regex": r'\b(?:DOB|D\.O\.B\.?|date\s+of\s+birth|birth\s*date|birth\s*day|b[\-\s]?day|born(?:\s+on)?|fecha\s+de\s+nacimiento)[\s:]*\d{1,2}[/\-\.]\d{1,2}[/\-\.]\d{2,4}\b',
                "description": "Date of Birth (labeled)",
                "valve": "block_phi",
                "severity": "high",
            },
            "phi_dob_iso": {
                "regex": r'\b(?:DOB|D\.O\.B\.?|date\s+of\s+birth|birth\s*date|birth\s*day|b[\-\s]?day|born(?:\s+on)?)[\s:]*\d{4}[/\-\.]\d{1,2}[/\-\.]\d{1,2}\b',
                "description": "Date of Birth - ISO format",
                "valve": "block_phi",
                "severity": "high",
            },
            "phi_dob_text_month": {
                "regex": r'\b(?:DOB|D\.O\.B\.?|date\s+of\s+birth|birth\s*date|birth\s*day|b[\-\s]?day|born(?:\s+on)?)[\s:]*(?:Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|Jun(?:e)?|Jul(?:y)?|Aug(?:ust)?|Sep(?:t(?:ember)?)?|Oct(?:ober)?|Nov(?:ember)?|Dec(?:ember)?)\.?\s+\d{1,2},?\s+\d{2,4}',
                "description": "Date of Birth - text month",
                "valve": "block_phi",
                "severity": "high",
            },
            "phi_dob_standalone": {
                "regex": r'\b(?:0?[1-9]|1[0-2])[/\-](?:0?[1-9]|[12]\d|3[01])[/\-](?:19|20)\d{2}\b',
                "description": "Date Pattern (MM/DD/YYYY)",
                "valve": "block_standalone_dates",
                "severity": "medium",
            },
            "phi_dob_standalone_iso": {
                "regex": r'\b(?:19|20)\d{2}[/\-](?:0?[1-9]|1[0-2])[/\-](?:0?[1-9]|[12]\d|3[01])\b',
                "description": "Date Pattern (YYYY-MM-DD)",
                "valve": "block_standalone_dates",
                "severity": "medium",
            },
            "phi_diagnosis": {
                "regex": r'\b(?:ICD[-\s]?(?:9|10)[-\s]?(?:CM|PCS)?[\s:#]*[A-Z]\d{2}(?:\.\d{1,4})?)\b',
                "description": "ICD Diagnosis Code",
                "valve": "block_phi",
                "severity": "medium",
            },
            "api_key": {
                "regex": r'\b(?:sk-[a-zA-Z0-9]{20,}|api[_\-]?key[\s=:]+["\']?[a-zA-Z0-9_\-]{16,})',
                "description": "API Key",
                "valve": "block_credentials",
                "severity": "critical",
            },
            "password_inline": {
                "regex": r'(?:password|passwd|pwd)[\s]*[=:]+[\s]*["\']?[^\s"\']{8,}',
                "description": "Inline Password",
                "valve": "block_credentials",
                "severity": "critical",
            },
            "connection_string": {
                "regex": r'(?:Server|Data Source|Host|Provider)=[^;\n]+;(?:.*?(?:Password|Pwd|User ID)=[^;\n]+)',
                "description": "Database Connection String",
                "valve": "block_credentials",
                "severity": "critical",
            },
            "aws_key": {
                "regex": r'\b(?:AKIA|ASIA)[A-Z0-9]{16}\b',
                "description": "AWS Access Key",
                "valve": "block_credentials",
                "severity": "critical",
            },
            "private_key": {
                "regex": r'-----BEGIN (?:RSA |EC |DSA )?PRIVATE KEY-----',
                "description": "Private Key",
                "valve": "block_credentials",
                "severity": "critical",
            },
            "bank_routing": {
                "regex": r'\b(?:routing|ABA)[\s#:]*\d{9}\b',
                "description": "Bank Routing Number",
                "valve": "block_bank_accounts",
                "severity": "critical",
            },
            "bank_account": {
                "regex": r'\b(?:account|acct)[\s#:]*\d{8,17}\b',
                "description": "Bank Account Number",
                "valve": "block_bank_accounts",
                "severity": "critical",
            },
        }

    # =========================================================================
    # Pattern matching
    # =========================================================================

    def _get_active_patterns(self) -> Dict[str, Dict]:
        """Return only patterns whose corresponding valve is enabled."""
        active = {}
        for name, pattern in self.patterns.items():
            valve_name = pattern.get("valve", "")
            if getattr(self.valves, valve_name, True):
                active[name] = pattern

        # Add custom patterns
        try:
            custom = json.loads(self.valves.custom_patterns)
            for name, regex in custom.items():
                active[f"custom_{name}"] = {
                    "regex": regex,
                    "description": f"Custom: {name}",
                    "valve": "enabled",
                    "severity": "high",
                }
        except (json.JSONDecodeError, TypeError):
            pass

        return active

    def _scan_text(self, text: str) -> List[Tuple[str, str, str]]:
        """
        Scan text for all active patterns.
        Returns list of (pattern_name, description, severity) for each match.
        """
        detections = []
        active_patterns = self._get_active_patterns()

        for name, pattern in active_patterns.items():
            if re.search(pattern["regex"], text, re.IGNORECASE):
                detections.append((
                    name,
                    pattern["description"],
                    pattern["severity"],
                ))

        return detections

    def _redact_text(self, text: str) -> str:
        """Replace all pattern matches with [REDACTED-TYPE] placeholders."""
        active_patterns = self._get_active_patterns()
        redacted = text

        for name, pattern in active_patterns.items():
            redacted = re.sub(
                pattern["regex"],
                f"[REDACTED-{name.upper()}]",
                redacted,
                flags=re.IGNORECASE,
            )

        return redacted

    # =========================================================================
    # File text extraction
    # =========================================================================

    def _resolve_file_path(self, file_id: str, file_name: str) -> Optional[str]:
        """Locate an uploaded file on disk by its ID."""
        # Open WebUI stores files as: <upload_dir>/<file_id>_<filename>
        # or just <upload_dir>/<file_id> depending on version
        candidates = glob_mod.glob(os.path.join(UPLOAD_DIR, f"{file_id}*"))
        if candidates:
            return candidates[0]

        # Fallback: search recursively (some versions nest by user)
        candidates = glob_mod.glob(os.path.join(UPLOAD_DIR, "**", f"{file_id}*"), recursive=True)
        if candidates:
            return candidates[0]

        return None

    def _get_file_extension(self, file_name: str) -> str:
        """Extract lowercase file extension from a filename."""
        if "." in file_name:
            return file_name.rsplit(".", 1)[-1].lower()
        return ""

    def _extract_text_from_file(self, file_path: str, file_name: str) -> str:
        """
        Extract text content from a file based on its extension.
        Returns extracted text or empty string if format is unsupported.
        """
        ext = self._get_file_extension(file_name)

        extractors = {
            "xlsx": self._extract_xlsx,
            "xlsm": self._extract_xlsx,
            "xls": self._extract_xls,
            "csv": self._extract_csv,
            "tsv": self._extract_csv,
            "pdf": self._extract_pdf,
            "docx": self._extract_docx,
            "pptx": self._extract_pptx,
        }

        # Plain text formats
        text_extensions = {
            "txt", "md", "json", "xml", "yaml", "yml", "log",
            "py", "js", "ts", "html", "css", "sql", "sh", "bat",
            "ini", "cfg", "conf", "env", "properties",
        }

        if ext in extractors:
            return extractors[ext](file_path)
        elif ext in text_extensions:
            return self._extract_plain_text(file_path)
        else:
            logger.info(f"DLP: Unsupported file format '.{ext}' for '{file_name}', skipping")
            return ""

    def _extract_xlsx(self, file_path: str) -> str:
        """Extract text from Excel .xlsx/.xlsm files."""
        try:
            import openpyxl
        except ImportError:
            logger.warning("DLP: openpyxl not installed, cannot scan .xlsx files")
            return ""

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    parts.append(row_text)
        wb.close()
        return "\n".join(parts)

    def _extract_xls(self, file_path: str) -> str:
        """Extract text from legacy Excel .xls files."""
        try:
            import xlrd
        except ImportError:
            logger.warning("DLP: xlrd not installed, cannot scan .xls files")
            return ""

        wb = xlrd.open_workbook(file_path)
        parts = []
        for sheet in wb.sheets():
            for row_idx in range(sheet.nrows):
                row_text = " ".join(
                    str(sheet.cell_value(row_idx, col_idx))
                    for col_idx in range(sheet.ncols)
                    if sheet.cell_value(row_idx, col_idx) not in (None, "")
                )
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)

    def _extract_csv(self, file_path: str) -> str:
        """Extract text from CSV/TSV files."""
        import csv
        parts = []
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                row_text = " ".join(cell for cell in row if cell)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)

    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF files."""
        try:
            from pypdf import PdfReader
        except ImportError:
            try:
                from PyPDF2 import PdfReader
            except ImportError:
                logger.warning("DLP: pypdf/PyPDF2 not installed, cannot scan .pdf files")
                return ""

        reader = PdfReader(file_path)
        parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                parts.append(text)
        return "\n".join(parts)

    def _extract_docx(self, file_path: str) -> str:
        """Extract text from Word .docx files."""
        try:
            import docx2txt
            return docx2txt.process(file_path) or ""
        except ImportError:
            pass

        try:
            from docx import Document
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            logger.warning("DLP: docx2txt/python-docx not installed, cannot scan .docx files")
            return ""

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint .pptx files."""
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("DLP: python-pptx not installed, cannot scan .pptx files")
            return ""

        prs = Presentation(file_path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            parts.append(text)
        return "\n".join(parts)

    def _extract_plain_text(self, file_path: str) -> str:
        """Extract text from plain text files."""
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    # =========================================================================
    # File scanning
    # =========================================================================

    def _scan_files(self, body: dict, user_id: str) -> List[Tuple[str, List[Tuple[str, str, str]]]]:
        """
        Scan all uploaded files in the request for sensitive data.

        Returns a list of (file_name, detections) tuples where detections is
        non-empty. Each detection is (pattern_name, description, severity).
        """
        files = body.get("files", [])
        if not files:
            return []

        max_bytes = self.valves.max_file_size_mb * 1024 * 1024
        flagged_files = []

        for file_ref in files:
            if not isinstance(file_ref, dict):
                continue

            file_type = file_ref.get("type", "")
            if file_type not in ("file", ""):
                continue

            file_id = file_ref.get("id", "")
            file_name = file_ref.get("name", "unknown")

            if not file_id:
                continue

            try:
                file_path = self._resolve_file_path(file_id, file_name)
                if not file_path:
                    logger.info(f"DLP: Could not locate file '{file_name}' (id={file_id}) on disk")
                    continue

                file_size = os.path.getsize(file_path)
                if file_size > max_bytes:
                    logger.warning(
                        f"DLP: File '{file_name}' ({file_size // (1024*1024)}MB) "
                        f"exceeds max scan size ({self.valves.max_file_size_mb}MB), skipping"
                    )
                    continue

                text = self._extract_text_from_file(file_path, file_name)
                if not text:
                    continue

                detections = self._scan_text(text)
                if detections:
                    flagged_files.append((file_name, detections))

                    if self.valves.log_detections:
                        detection_summary = ", ".join(
                            f"{d[1]} ({d[2]})" for d in detections
                        )
                        logger.warning(
                            f"DLP: Sensitive data in file '{file_name}' "
                            f"from user={user_id}: {detection_summary}"
                        )

            except Exception as e:
                logger.error(f"DLP: Error scanning file '{file_name}': {e}")
                continue

        return flagged_files

    # =========================================================================
    # Pipeline hooks
    # =========================================================================

    async def inlet(self, body: dict, __user__: Optional[dict] = None) -> dict:
        """
        Pre-processing filter: runs on every message before it reaches the LLM.

        Scans both message text and uploaded files for sensitive data.
        In 'block' mode: rejects the entire request if sensitive data is found.
        In 'redact' mode: replaces sensitive text and strips flagged files.
        """
        if not self.valves.enabled:
            return body

        messages = body.get("messages", [])
        user_id = __user__.get("id", "unknown") if __user__ else "unknown"

        # --- Sanitize empty content blocks (Claude API rejects empty text) ---
        for message in messages:
            content = message.get("content", "")
            if isinstance(content, list):
                for item in content:
                    if (isinstance(item, dict)
                            and item.get("type") == "text"
                            and not item.get("text", "").strip()):
                        item["text"] = "Please analyze the attached file."
            elif isinstance(content, str) and not content.strip():
                message["content"] = "Please analyze the attached file."

        # --- Scan message text ---
        message_detections = []
        for i, message in enumerate(messages):
            if message.get("role") != "user":
                continue

            content = message.get("content", "")
            if isinstance(content, list):
                text_parts = [
                    item.get("text", "")
                    for item in content
                    if isinstance(item, dict) and item.get("type") == "text"
                ]
                text_to_scan = " ".join(text_parts)
            else:
                text_to_scan = str(content)

            detections = self._scan_text(text_to_scan)
            if detections:
                message_detections.extend(detections)

                if self.valves.log_detections:
                    detection_summary = ", ".join(
                        f"{d[1]} ({d[2]})" for d in detections
                    )
                    logger.warning(
                        f"DLP: Detected sensitive data from user={user_id}: "
                        f"{detection_summary}"
                    )

        # --- Scan uploaded files ---
        file_detections = []
        if self.valves.scan_file_uploads:
            file_detections = self._scan_files(body, user_id)

        # --- No detections: pass through ---
        if not message_detections and not file_detections:
            return body

        # --- Handle detections ---
        if self.valves.mode == "block":
            error_parts = []

            if message_detections:
                types_found = ", ".join(set(d[1] for d in message_detections))
                error_parts.append(
                    f"Your message text contains sensitive information "
                    f"({types_found})."
                )

            if file_detections:
                file_lines = []
                for fname, dets in file_detections:
                    types_found = ", ".join(set(d[1] for d in dets))
                    file_lines.append(f"  - **{fname}**: {types_found}")
                error_parts.append(
                    "Uploaded files contain sensitive information:\n"
                    + "\n".join(file_lines)
                )

            raise Exception(
                f"**DLP Filter Blocked This Message**\n\n"
                + "\n\n".join(error_parts)
                + "\n\nFor security and compliance, this message has been "
                f"blocked from being sent to the AI service.\n\n"
                f"Please remove the sensitive data and try again."
            )

        elif self.valves.mode == "redact":
            # Redact message text
            if message_detections:
                for message in messages:
                    if message.get("role") != "user":
                        continue
                    content = message.get("content", "")
                    if isinstance(content, list):
                        for item in content:
                            if isinstance(item, dict) and item.get("type") == "text":
                                item["text"] = self._redact_text(item["text"])
                    else:
                        message["content"] = self._redact_text(str(content))

            # Strip flagged files from the request
            if file_detections:
                flagged_names = {fname for fname, _ in file_detections}
                original_files = body.get("files", [])
                clean_files = [
                    f for f in original_files
                    if f.get("name") not in flagged_names
                ]
                body["files"] = clean_files

                # Also strip from metadata if present
                metadata = body.get("metadata", {})
                if "files" in metadata:
                    metadata["files"] = [
                        f for f in metadata["files"]
                        if f.get("name") not in flagged_names
                    ]

                # Notify the user which files were removed
                stripped_summary = ", ".join(flagged_names)
                body["messages"].append({
                    "role": "system",
                    "content": (
                        f"[DLP Notice] The following files were removed because they "
                        f"contain sensitive data: {stripped_summary}. "
                        f"The AI will not see these files."
                    ),
                })

        return body

    async def outlet(self, body: dict, __user__: Optional[dict] = None) -> dict:
        """
        Post-processing filter: runs on every response from the LLM.
        Scans assistant responses for any echoed-back sensitive data.
        """
        if not self.valves.enabled:
            return body

        messages = body.get("messages", [])

        for message in messages:
            if message.get("role") != "assistant":
                continue

            content = message.get("content", "")
            if isinstance(content, str):
                detections = self._scan_text(content)
                if detections:
                    message["content"] = self._redact_text(content)
                    if self.valves.log_detections:
                        logger.warning(
                            f"DLP: Redacted sensitive data in assistant response"
                        )

        return body


# Open WebUI v0.8+ expects a class named "Filter" or "Function" for filter modules
Filter = Pipeline
Function = Pipeline
